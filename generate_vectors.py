from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def create_vector_slides():
    prs = Presentation()
    # 使用空白布局
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 统一定义字体和颜色（原图偏向 Serif 衬线体）
    FONT_NAME = 'Times New Roman'
    BLACK = RGBColor(0, 0, 0)

    # --- 辅助函数区 ---
    def add_bounding_box(top_y):
        """添加外层圆角矩形"""
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top_y), Inches(9.0), Inches(2.8))
        rect.fill.background()  # 透明背景
        rect.line.color.rgb = RGBColor(100, 100, 100)  # 灰色边线
        rect.line.width = Pt(1)
        return rect

    def add_label(left, top, width, height, text_runs):
        """添加左侧的多段格式文本（支持斜体和下标）"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT

        for t, is_italic, is_sub, is_bold in text_runs:
            run = p.add_run()
            run.text = t
            run.font.size = Pt(16)
            run.font.name = FONT_NAME
            run.font.color.rgb = BLACK
            if is_italic: run.font.italic = True
            if is_sub: run.font.subscript = True
            if is_bold: run.font.bold = True
        return txBox

    def add_center_text(text, left, top, width, height, font_size, is_italic=False):
        """添加顶部居中的文本（位置、任务等）"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = FONT_NAME
        run.font.color.rgb = BLACK
        run.font.italic = is_italic
        return txBox

    def draw_grid_boxes(start_x, start_y, box_w, box_h, values):
        """绘制连续的方格和内部数值"""
        for i, val in enumerate(values):
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x + i * box_w), Inches(start_y),
                                          Inches(box_w), Inches(box_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
            rect.line.color.rgb = BLACK
            rect.line.width = Pt(1)

            tf = rect.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(18)
            run.font.name = FONT_NAME
            run.font.color.rgb = BLACK

    def add_caption(left, top, width, height, fig_num, caption_text):
        """添加左下角的图注"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        r1 = p.add_run()
        r1.text = fig_num
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.name = FONT_NAME
        r1.font.color.rgb = BLACK

        r2 = p.add_run()
        r2.text = " " + caption_text
        r2.font.size = Pt(14)
        r2.font.name = FONT_NAME
        r2.font.color.rgb = BLACK

    # ==========================================
    # 绘制 Figure 3 (位于幻灯片上半部分)
    # ==========================================
    top_fig3 = 0.5
    add_bounding_box(top_fig3)

    # 左侧标签: (文本, 是否斜体, 是否下标, 是否加粗)
    add_label(0.5, top_fig3 + 0.3, 3.5, 0.4, [("Position: Priority ", False, False, False), ("u", True, False, False)])
    add_label(0.5, top_fig3 + 0.8, 3.5, 0.4, [("Task Indicated", False, False, False)])
    add_label(0.5, top_fig3 + 1.4, 3.5, 0.4,
              [("Transport sequence v", False, False, False), ("1", False, True, False), ("(u)", False, False, False)])

    # 右侧排版坐标
    start_x = 4.2
    box_w = 0.8

    # 顶部位置数字
    for i in range(5):
        add_center_text(str(i + 1), start_x + i * box_w, top_fig3 + 0.3, box_w, 0.4, 18)

    # 中间任务名称 (注意原图中 Task 序号打乱)
    tasks_fig3 = ["Task3", "Task2", "Task1", "Task4", "Task5"]
    for i in range(5):
        add_center_text(tasks_fig3[i], start_x + i * box_w, top_fig3 + 0.8, box_w, 0.4, 16, is_italic=True)

    # 底部方格矩阵
    draw_grid_boxes(start_x, top_fig3 + 1.35, box_w, 0.6, ["3", "2", "1", "4", "5"])

    # 左下角图注
    add_caption(0.7, top_fig3 + 2.2, 8.0, 0.4, "Figure 3", "Illustration of the transport sequence vector")

    # ==========================================
    # 绘制 Figure 4 (位于幻灯片下半部分)
    # ==========================================
    top_fig4 = 3.8
    add_bounding_box(top_fig4)

    # 左侧标签
    add_label(0.5, top_fig4 + 0.3, 3.5, 0.4, [("Position: ", False, False, False), ("v", True, False, False)])
    add_label(0.5, top_fig4 + 0.8, 3.5, 0.4, [("Task Indicated", False, False, False)])
    add_label(0.5, top_fig4 + 1.4, 3.5, 0.4,
              [("Transporter Assignment v", False, False, False), ("1", False, True, False),
               ("(v)", False, False, False)])

    # 顶部位置数字
    for i in range(5):
        add_center_text(str(i + 1), start_x + i * box_w, top_fig4 + 0.3, box_w, 0.4, 18)

    # 中间任务名称 (顺排)
    tasks_fig4 = ["Task1", "Task2", "Task3", "Task4", "Task5"]
    for i in range(5):
        add_center_text(tasks_fig4[i], start_x + i * box_w, top_fig4 + 0.8, box_w, 0.4, 16, is_italic=True)

    # 底部方格矩阵
    draw_grid_boxes(start_x, top_fig4 + 1.35, box_w, 0.6, ["2", "2", "1", "2", "1"])

    # 左下角图注
    add_caption(0.7, top_fig4 + 2.2, 8.0, 0.4, "Figure 4", "Illustration of the transporter assignment vector")

    # 保存幻灯片
    output_filename = "vector_illustrations.pptx"
    prs.save(output_filename)
    print(f"PPT 生成成功！已保存为: {output_filename}")


if __name__ == '__main__':
    create_vector_slides()