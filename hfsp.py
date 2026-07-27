from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 定义原图中的主题颜色
def create_hfsp_diagram():
    # 创建一个新的PPT，选择一张空白幻灯片 (布局索引为6)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 定义原图中的主题颜色
    color_blue = RGBColor(180, 210, 235)  # Stage 1
    color_green = RGBColor(168, 208, 141)  # Stage 2
    color_grey = RGBColor(217, 217, 217)  # Stage C
    color_white = RGBColor(255, 255, 255)  # 背景/纯白

    # 封装一个快捷创建带有文本的矩形框的函数
    def add_box(x, y, w, h, text, fill_color, border=True):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        # 填充颜色
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

        # 边框设置
        if not border:
            shape.line.fill.background()  # 隐藏边框
        else:
            shape.line.color.rgb = RGBColor(0, 0, 0)

        # 文本设置
        shape.text = text
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.name = 'Arial'
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.bold = True
        return shape

    # 1. 左侧入口：Job Queue (仓库或装卸区)
    add_box(0.5, 2.5, 0.8, 2.5, "JOB\nQUEUE", color_white)

    # 2. 顶部阶段标题 (无边框)
    add_box(2.0, 0.5, 1.5, 0.5, "STAGE 1", color_white, border=False)
    add_box(4.5, 0.5, 1.5, 0.5, "STAGE 2", color_white, border=False)
    add_box(7.0, 0.5, 1.5, 0.5, "STAGE C", color_white, border=False)

    # 3. Stage 1 机器矩阵 (蓝色)
    add_box(2.0, 1.5, 1.5, 1.0, "MACHINE\nM 1,1", color_blue)
    add_box(2.0, 3.0, 1.5, 1.0, "MACHINE\nM 1,2", color_blue)
    add_box(2.0, 4.2, 1.5, 0.5, "...", color_white, border=False)  # 省略号
    add_box(2.0, 5.0, 1.5, 1.0, "MACHINE\nM 1,x", color_blue)

    # 4. Stage 2 机器矩阵 (绿色)
    add_box(4.5, 1.5, 1.5, 1.0, "MACHINE\nM 2,1", color_green)
    add_box(4.5, 3.0, 1.5, 1.0, "MACHINE\nM 2,2", color_green)
    add_box(4.5, 4.2, 1.5, 0.5, "...", color_white, border=False)
    add_box(4.5, 5.0, 1.5, 1.0, "MACHINE\nM 2,y", color_green)

    # 5. Stage C 机器矩阵 (灰色)
    add_box(7.0, 1.5, 1.5, 1.0, "MACHINE\nM C,1", color_grey)
    add_box(7.0, 3.0, 1.5, 1.0, "MACHINE\nM C,2", color_grey)
    add_box(7.0, 4.2, 1.5, 0.5, "...", color_white, border=False)
    add_box(7.0, 5.0, 1.5, 1.0, "MACHINE\nM C,z", color_grey)

    # 6. 右侧出口：Completed Workpieces
    add_box(9.0, 2.5, 0.8, 2.5, "COMPLETED\nWORKPIECES", color_white)

    # 保存 PPT 文件
    prs.save('Editable_Production_Topology.pptx')
    print("生成成功！请在当前目录下找到 Editable_Production_Topology.pptx 文件。")


if __name__ == '__main__':
    create_hfsp_diagram()