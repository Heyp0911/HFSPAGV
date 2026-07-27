from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn


def add_arrow(slide, start_x, start_y, end_x, end_y):
    """绘制带有箭头的实线"""
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(1.5)

    # 通过更稳定的底层 XML XPath 添加箭头端点
    try:
        ln = connector._element.xpath('.//a:ln')[0]
        tailEnd = ln.find(qn('a:tailEnd'))
        if tailEnd is None:
            tailEnd = OxmlElement('a:tailEnd')
            ln.append(tailEnd)
        tailEnd.set('type', 'triangle')
        # 可选：设置箭头大小
        tailEnd.set('w', 'med')
        tailEnd.set('len', 'med')
    except Exception as e:
        print(f"警告: 添加箭头失败 - {e}")

    return connector


def add_line(slide, start_x, start_y, end_x, end_y):
    """绘制不带箭头的实线"""
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(1.5)
    return connector


def add_textbox(slide, text, x, y, w, h, bold=False):
    """添加透明底色的纯文本框"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.bold = bold
    return txBox


def add_box(slide, text, x, y, w, h, fill_rgb, line_rgb, bold=False):
    """添加带有背景色和边框的处理节点"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # 轻微圆角
    shape.adjustments[0] = 0.1

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.2)

    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.bold = bold
    return shape


def add_background(slide, title, x, y, w, h, fill_rgb, line_rgb, title_color):
    """添加虚线框背景组"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.5)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # 标题文字
    txBox = slide.shapes.add_textbox(x, y - Inches(0.1), w, Inches(0.5))
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    # 使用 16:9 宽屏布局 (13.33 x 7.5 英寸)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ================= 1. 绘制背景层 (先画在底层) =================
    # Parallel Preparation Stage
    add_background(slide, "Parallel Preparation Stage",
                   Inches(0.4), Inches(1.1), Inches(4.3), Inches(3.4),
                   RGBColor(245, 245, 245), RGBColor(150, 150, 150), RGBColor(51, 51, 51))

    # Main Assembly and Testing Stage
    add_background(slide, "Main Assembly and Testing Stage",
                   Inches(5.0), Inches(1.1), Inches(7.9), Inches(3.4),
                   RGBColor(253, 242, 235), RGBColor(204, 102, 0), RGBColor(204, 102, 0))

    # ================= 2. 绘制连接线 (放在节点后方避免穿模) =================
    # 电池分支持续直线
    add_arrow(slide, Inches(1.7), Inches(2.1), Inches(2.5), Inches(2.1))
    add_arrow(slide, Inches(4.5), Inches(2.1), Inches(5.3), Inches(2.1))

    # 柜体分支持续直线与折线汇集
    add_arrow(slide, Inches(1.7), Inches(3.9), Inches(2.5), Inches(3.9))
    # Op2 -> Op3 (手动绘制直角折线)
    add_line(slide, Inches(4.5), Inches(3.9), Inches(4.9), Inches(3.9))
    add_line(slide, Inches(4.9), Inches(3.9), Inches(4.9), Inches(2.1))
    add_arrow(slide, Inches(4.9), Inches(2.1), Inches(5.3), Inches(2.1))

    # 后续串行工序连接线
    add_arrow(slide, Inches(7.1), Inches(2.1), Inches(7.6), Inches(2.1))  # Op3 -> Op4
    add_arrow(slide, Inches(9.1), Inches(2.1), Inches(9.6), Inches(2.1))  # Op4 -> Op5
    add_arrow(slide, Inches(10.6), Inches(2.4), Inches(10.6), Inches(3.6))  # Op5 -> Op6 (向下)
    add_arrow(slide, Inches(9.6), Inches(3.9), Inches(9.1), Inches(3.9))  # Op6 -> Op7 (向左)
    add_arrow(slide, Inches(7.6), Inches(3.9), Inches(7.1), Inches(3.9))  # Op7 -> Finish (向左)

    # ================= 3. 绘制节点模块 =================
    # 输入原料文本
    add_textbox(slide, "Battery Cells", Inches(0.5), Inches(1.8), Inches(1.2), Inches(0.6), bold=True)
    add_textbox(slide, "Cabinet Base", Inches(0.5), Inches(3.6), Inches(1.2), Inches(0.6), bold=True)
    add_textbox(slide, "Finished Cabinet", Inches(5.3), Inches(3.6), Inches(1.8), Inches(0.6), bold=True)

    # 并行阶段节点
    add_box(slide, "Op1: Module and\nPack Assembly",
            Inches(2.5), Inches(1.8), Inches(2.0), Inches(0.6),
            RGBColor(230, 240, 250), RGBColor(0, 51, 153))

    add_box(slide, "Op2: Cabinet Frame\nPre-assembly",
            Inches(2.5), Inches(3.6), Inches(2.0), Inches(0.6),
            RGBColor(230, 250, 230), RGBColor(0, 102, 0))

    # 串行合并阶段节点
    add_box(slide, "Op3: Pack-into-\nCabinet",
            Inches(5.3), Inches(1.8), Inches(1.8), Inches(0.6),
            RGBColor(255, 255, 255), RGBColor(0, 0, 0))

    add_box(slide, "Op4: Pack\nFastening",
            Inches(7.6), Inches(1.8), Inches(1.5), Inches(0.6),
            RGBColor(255, 255, 255), RGBColor(0, 0, 0))

    add_box(slide, "Op5: System Cabling\nIntegration",
            Inches(9.6), Inches(1.8), Inches(2.0), Inches(0.6),
            RGBColor(255, 255, 255), RGBColor(0, 0, 0))

    add_box(slide, "Op6: Insulation and\nAirtightness Test",
            Inches(9.6), Inches(3.6), Inches(2.0), Inches(0.6),
            RGBColor(255, 255, 255), RGBColor(0, 0, 0))

    add_box(slide, "Op7: Final\nFactory Test",
            Inches(7.6), Inches(3.6), Inches(1.5), Inches(0.6),
            RGBColor(255, 255, 255), RGBColor(0, 0, 0))

    # 保存文件
    prs.save("Figure2_Assembly_Flowchart.pptx")
    print("成功生成 PPT 流程图：Figure2_Assembly_Flowchart.pptx")


if __name__ == '__main__':
    main()