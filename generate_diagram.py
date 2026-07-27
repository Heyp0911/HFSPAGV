from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE


def create_task_pool_slide():
    prs = Presentation()
    # 使用空白布局
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # --- 辅助函数：添加灰色背景的标题头 ---
    def add_header(text, left, top, width, height):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(217, 217, 217)  # 灰色背景
        rect.line.fill.background()  # 无边框
        tf = rect.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.size = Pt(14)
        return rect

    # --- 辅助函数：添加普通任务框（支持部分斜体） ---
    def add_task_box(text, left, top, width, height, border_color, dashed=False):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.line.color.rgb = border_color
        if dashed:
            rect.line.dash_style = MSO_LINE.DASH

        tf = rect.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # 处理斜体 (分离 "TaskX" 和 括号内容)
        parts = text.split(" ", 1)

        r1 = p.add_run()
        r1.text = parts[0] + " "
        r1.font.italic = True
        r1.font.color.rgb = RGBColor(0, 0, 0)
        r1.font.size = Pt(12)

        if len(parts) > 1:
            r2 = p.add_run()
            r2.text = parts[1]
            r2.font.color.rgb = RGBColor(0, 0, 0)
            r2.font.size = Pt(12)
        return rect

    # 1. 绘制顶部标题
    add_header("New Task", Inches(0.8), Inches(1), Inches(1.8), Inches(0.4))
    add_header("Task Pool", Inches(4.1), Inches(1), Inches(1.8), Inches(0.4))
    add_header("Old Task", Inches(7.4), Inches(1), Inches(1.8), Inches(0.4))

    # 2. 绘制 New Task 的多色文本框
    new_task_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(2.5), Inches(2.9), Inches(0.4))
    new_task_rect.fill.solid()
    new_task_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    new_task_rect.line.color.rgb = RGBColor(0, 0, 0)

    tf = new_task_rect.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    r1 = p.add_run()
    r1.text = "Task2 "
    r1.font.italic = True
    r1.font.size = Pt(12)
    r1.font.color.rgb = RGBColor(0, 0, 0)

    r2 = p.add_run()
    r2.text = "(job2, "
    r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor(0, 0, 0)

    r3 = p.add_run()
    r3.text = "stage2, 6, machine2"
    r3.font.color.rgb = RGBColor(255, 0, 0)  # 红色字体
    r3.font.size = Pt(12)

    r4 = p.add_run()
    r4.text = ")"
    r4.font.size = Pt(12)
    r4.font.color.rgb = RGBColor(0, 0, 0)

    # 3. 绘制 "add" 箭头
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(2.5), Inches(0.5), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(255, 0, 0)
    arrow1.line.fill.background()

    tx_add = slide.shapes.add_textbox(Inches(3.1), Inches(2.8), Inches(0.7), Inches(0.3))
    tx_add.text_frame.text = "add"
    tx_add.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tx_add.text_frame.paragraphs[0].font.bold = True
    tx_add.text_frame.paragraphs[0].font.size = Pt(11)

    # 4. 绘制中间的 Task Pool 蓝色大圆角矩形
    pool = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(1.6), Inches(2.6), Inches(2.4))
    pool.fill.background()
    pool.line.color.rgb = RGBColor(68, 114, 196)  # 蓝色边线
    pool.line.width = Pt(1.5)

    # 5. 填充 Pool 内部的任务
    add_task_box("Task1 (job1, stage1, 0, warehouse)", Inches(3.9), Inches(1.8), Inches(2.4), Inches(0.4),
                 RGBColor(0, 0, 0))
    add_task_box("Task2 (job2, stage1, 0, warehouse)", Inches(3.9), Inches(2.6), Inches(2.4), Inches(0.4),
                 RGBColor(255, 0, 0), dashed=True)
    add_task_box("Task3 (job3, stage1, 0, warehouse)", Inches(3.9), Inches(3.4), Inches(2.4), Inches(0.4),
                 RGBColor(0, 0, 0))

    # 6. 绘制 "remove" 箭头
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(2.5), Inches(0.5), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(255, 0, 0)
    arrow2.line.fill.background()

    tx_fin = slide.shapes.add_textbox(Inches(6.4), Inches(2.2), Inches(0.7), Inches(0.3))
    tx_fin.text_frame.text = "finished"
    tx_fin.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tx_fin.text_frame.paragraphs[0].font.bold = True
    tx_fin.text_frame.paragraphs[0].font.size = Pt(11)

    tx_rem = slide.shapes.add_textbox(Inches(6.4), Inches(2.8), Inches(0.7), Inches(0.3))
    tx_rem.text_frame.text = "remove"
    tx_rem.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tx_rem.text_frame.paragraphs[0].font.bold = True
    tx_rem.text_frame.paragraphs[0].font.size = Pt(11)

    # 7. 绘制 Old Task 盒子
    add_task_box("Task2 (job2, stage1, 0, warehouse)", Inches(7.1), Inches(2.5), Inches(2.6), Inches(0.4),
                 RGBColor(128, 128, 128), dashed=True)

    # 8. 左下角图注
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(5), Inches(0.4))
    p_cap = caption.text_frame.paragraphs[0]
    run_cap1 = p_cap.add_run()
    run_cap1.text = "Figure 2 "
    run_cap1.font.bold = True
    run_cap1.font.size = Pt(12)

    run_cap2 = p_cap.add_run()
    run_cap2.text = "Simple description of the task pool"
    run_cap2.font.size = Pt(12)

    # 保存幻灯片
    output_filename = "task_pool_diagram.pptx"
    prs.save(output_filename)
    print(f"PPT 成功生成，已保存为: {output_filename}")


if __name__ == '__main__':
    create_task_pool_slide()